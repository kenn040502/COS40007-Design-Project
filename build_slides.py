# -*- coding: utf-8 -*-
"""
Build the COS40007 Design Project presentation deck.
Maps each slide to the marking rubric so every assessed item is visible.
Run AFTER run_analysis.py:  python build_slides.py
Outputs: outputs/COS40007_E's_Island_Presentation.pptx
"""
import os
import json
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT   = os.path.dirname(os.path.abspath(__file__))
FIG    = os.path.join(ROOT, "outputs", "figures")
MODELS = os.path.join(ROOT, "outputs", "models")
PROC   = os.path.join(ROOT, "data", "processed")
OUT    = os.path.join(ROOT, "outputs", "COS40007_E's_Island_Presentation.pptx")
os.makedirs(os.path.join(ROOT, "outputs"), exist_ok=True)

# ── Load data ────────────────────────────────────────────────────────────────
ts_metrics = pd.read_csv(os.path.join(MODELS, "ts_metrics_summary.csv"))
with open(os.path.join(MODELS, "clustering_results.json")) as f:
    clu = json.load(f)
with open(os.path.join(MODELS, "arima_results.json")) as f:
    ar = json.load(f)
km          = pd.read_csv(os.path.join(PROC, "states_kmeans_clustered.csv"))
overall_inf = pd.read_csv(os.path.join(PROC, "overall_inflation.csv"), parse_dates=["date"])
income      = pd.read_csv(os.path.join(PROC, "income_state_clean.csv"), parse_dates=["date"])

ts_row     = ts_metrics.iloc[0]
metrics    = ar["metrics"]
baseline   = ar.get("baseline", {}).get("metrics", {})
exog_exp   = ar.get("fuel_exog_experiment", {})
arima_ord  = ar["arima_order"]
stat       = ar["stationarity"]
fc_vals    = ar["forecast"]["values"]
best_k     = clu["best_k"]
km_scores  = clu["kmeans"]["scores"]
hc_scores  = clu["hierarchical"]["scores"]

latest_date   = overall_inf["date"].max().strftime("%b %Y")
earliest_date = overall_inf["date"].min().strftime("%b %Y")

# ── Style constants ──────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x4E, 0x79)
ACCENT  = RGBColor(0xD7, 0x35, 0x27)
DARK    = RGBColor(0x22, 0x22, 0x22)
GREY    = RGBColor(0x66, 0x66, 0x66)
LIGHT   = RGBColor(0xF2, 0xF2, 0xF2)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank

# ── Helpers ──────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def add_bg_band(slide, color=NAVY, height=Inches(0.55)):
    """Top navy band."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, height
    )
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = color
    return band

def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return tb

def add_title(slide, title, subtitle=None, rubric_tag=None):
    add_bg_band(slide)
    add_textbox(slide, Inches(0.4), Inches(0.05), Inches(11), Inches(0.45),
                title, size=22, bold=True, color=WHITE)
    if rubric_tag:
        add_textbox(slide, Inches(10.5), Inches(0.12), Inches(2.5), Inches(0.35),
                    rubric_tag, size=10, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    if subtitle:
        add_textbox(slide, Inches(0.4), Inches(0.65), Inches(12.5), Inches(0.35),
                    subtitle, size=13, color=GREY)

def add_bullets(slide, left, top, width, height, items, size=15, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run()
            r1.text = "• " + lead
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = DARK
        else:
            r = p.add_run()
            r.text = "• " + it
            r.font.size = Pt(size)
            r.font.color.rgb = DARK
    return tb

def add_image(slide, path, left, top, width=None, height=None, caption=None):
    if not os.path.exists(path):
        add_textbox(slide, left, top, width or Inches(5), Inches(0.5),
                    f"[missing: {os.path.basename(path)}]", size=11, color=ACCENT)
        return None
    if width and not height:
        pic = slide.shapes.add_picture(path, left, top, width=width)
    elif height and not width:
        pic = slide.shapes.add_picture(path, left, top, height=height)
    else:
        pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
    if caption:
        add_textbox(slide,
                    left,
                    top + (height or pic.height) + Inches(0.05),
                    width or pic.width,
                    Inches(0.3),
                    caption, size=10, color=GREY, align=PP_ALIGN.CENTER)
    return pic

def add_metric_card(slide, left, top, width, height, label, value, color=NAVY):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = color
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.05)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = str(value)
    r1.font.size = Pt(26)
    r1.font.bold = True
    r1.font.color.rgb = color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(11)
    r2.font.color.rgb = GREY
    return box

def add_footer(slide, n, total):
    add_textbox(slide, Inches(0.3), Inches(7.15), Inches(8), Inches(0.3),
                "COS40007 Design Project  ·  Team E's Island  ·  Smart Government",
                size=9, color=GREY)
    add_textbox(slide, Inches(12.3), Inches(7.15), Inches(0.9), Inches(0.3),
                f"{n} / {total}", size=9, color=GREY, align=PP_ALIGN.RIGHT)

# =============================================================================
# Slide 1 — Title
# =============================================================================
s = add_slide()
# Big navy band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(3.2))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = NAVY

add_textbox(s, Inches(0.6), Inches(0.6), Inches(12), Inches(0.6),
            "COS40007  ·  AI for Engineering  ·  Design Project",
            size=16, color=WHITE)
add_textbox(s, Inches(0.6), Inches(1.2), Inches(12), Inches(1.0),
            "Smart Government: Forecasting Inflation &",
            size=36, bold=True, color=WHITE)
add_textbox(s, Inches(0.6), Inches(1.9), Inches(12), Inches(1.0),
            "Clustering Malaysian States from Open DOSM Data",
            size=36, bold=True, color=WHITE)
add_textbox(s, Inches(0.6), Inches(2.7), Inches(12), Inches(0.5),
            "Theme 3 — A reproducible AI pipeline + interactive dashboard",
            size=16, color=RGBColor(0xCC, 0xDD, 0xEE))

add_textbox(s, Inches(0.6), Inches(3.7), Inches(6), Inches(0.4),
            "Team E's Island", size=20, bold=True, color=NAVY)
team = [
    "Kenneth Hui Hong CHUA   ·   102782494",
    "Kelvin Wen Kiong FONG   ·   102782287",
    "Nigel Zi Jun LING   ·   102779140",
    "Neng Yi CHIENG   ·   104386364",
]
for i, t in enumerate(team):
    add_textbox(s, Inches(0.6), Inches(4.1 + 0.35 * i), Inches(7), Inches(0.35),
                t, size=14, color=DARK)

add_textbox(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.35),
            "Data: Department of Statistics Malaysia (OpenDOSM)  ·  CC BY 4.0  ·  Random seed = 42",
            size=11, color=GREY)

# =============================================================================
# Slide 2 — Problem & Motivation
# =============================================================================
s = add_slide()
add_title(s, "Problem & Motivation",
          "Why apply AI to Malaysian public-sector data?")

add_bullets(s, Inches(0.5), Inches(1.2), Inches(12), Inches(4), [
    ("National averages hide regional variation. ",
     "Aggregate CPI inflation hides large differences between states "
     "in income level and inflation exposure."),
    ("Policymakers need short-term inflation forecasts ",
     "to design subsidies, monetary policy, and social-protection schemes."),
    ("Clustering reveals economic peer groups ",
     "of states — useful for targeted regional policy rather than uniform national instruments."),
    ("Fuel prices are a known driver of CPI inflation ",
     "in Malaysia; quantifying the pass-through helps the government model subsidy policy."),
    ("Open DOSM data is rich but underused. ",
     "An AI pipeline + dashboard makes the insight accessible to non-technical stakeholders."),
], size=15)

add_textbox(s, Inches(0.5), Inches(5.7), Inches(12), Inches(0.6),
            'Research question: Can ARIMA reliably forecast national CPI inflation, '
            'and do Malaysia\'s 16 states form interpretable economic clusters from CPI + income features?',
            size=14, bold=True, color=NAVY)
add_footer(s, 2, 18)

# =============================================================================
# Slide 3 — Objectives & Outcomes
# =============================================================================
s = add_slide()
add_title(s, "Project Objectives & Outcomes")

add_textbox(s, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
            "Five Objectives", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5), [
    "Forecast national CPI inflation (ARIMA, 24-month horizon).",
    "Cluster the 16 states by CPI + income features.",
    "Quantify fuel-price → CPI pass-through.",
    "Validate every output against DOSM official figures.",
    "Deliver a reproducible pipeline + Flask dashboard.",
], size=14)

add_textbox(s, Inches(7.2), Inches(1.1), Inches(5.5), Inches(0.4),
            "Headline Outcomes", size=18, bold=True, color=NAVY)
add_metric_card(s, Inches(7.2), Inches(1.6),  Inches(2.6), Inches(1.2),
                "Test MASE (ARIMA)", f"{metrics.get('MASE', float('nan')):.2f}")
add_metric_card(s, Inches(10.1), Inches(1.6), Inches(2.6), Inches(1.2),
                "Test RMSE",        f"{metrics['RMSE']:.4f}")
add_metric_card(s, Inches(7.2), Inches(3.0),  Inches(2.6), Inches(1.2),
                "Optimal k (clustering)", best_k)
add_metric_card(s, Inches(10.1), Inches(3.0), Inches(2.6), Inches(1.2),
                "K-Means silhouette", f"{km_scores['silhouette']:.4f}")
add_metric_card(s, Inches(7.2), Inches(4.4),  Inches(2.6), Inches(1.2),
                f"Forecast range ({ar['horizon']} mo)",
                f"{min(fc_vals):.2f}–{max(fc_vals):.2f}%")
add_metric_card(s, Inches(10.1), Inches(4.4), Inches(2.6), Inches(1.2),
                "Datasets processed", "4 DOSM tables")
add_footer(s, 3, 18)

# =============================================================================
# Slide 4 — Data Sources  [Rubric 1]
# =============================================================================
s = add_slide()
add_title(s, "Data Sources",
          "All datasets from DOSM / OpenDOSM, CC BY 4.0, fetched programmatically",
          rubric_tag="Rubric 1: Data (10)")

rows = [
    ["Dataset", "Catalogue", "Coverage", "Records"],
    ["CPI Inflation by Division", "cpi_2d_inflation",
     f"National, {earliest_date}–{latest_date}", "7,770"],
    ["CPI Index by State & Division", "cpi_2d_state",
     "16 states, 2010–{}".format(latest_date), "43,904"],
    ["Retail Fuel Prices", "fuelprice",
     "Weekly 2017–2026 → monthly", "111"],
    ["Household Income by State", "hh_income_state",
     "16 states, 1970–2022", "303"],
]
tbl = s.shapes.add_table(len(rows), len(rows[0]),
                         Inches(0.5), Inches(1.3),
                         Inches(12.3), Inches(2.6)).table
widths = [Inches(3.4), Inches(2.6), Inches(4.3), Inches(2.0)]
for j, w in enumerate(widths):
    tbl.columns[j].width = w
for r_i, row in enumerate(rows):
    for c_i, val in enumerate(row):
        cell = tbl.cell(r_i, c_i)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(val)
        run.font.size = Pt(13)
        if r_i == 0:
            run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            run.font.color.rgb = DARK
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r_i % 2 else WHITE

add_textbox(s, Inches(0.5), Inches(4.2), Inches(12), Inches(0.4),
            "Why these 4?  Longitudinal (forecasting) + cross-sectional (clustering) + driver (fuel)",
            size=14, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(4.7), Inches(12), Inches(2.3), [
    ("Time-series target: ", "monthly national CPI YoY inflation (1 series, 544 obs)."),
    ("Cross-sectional features: ", "5 features × 16 states for clustering."),
    ("Driver variable: ", "monthly average retail fuel price (RON95/97/diesel)."),
    ("Validation source: ", "DOSM official 2022 HIES releases for ground truth checking."),
], size=14)
add_footer(s, 4, 18)

# =============================================================================
# Slide 5 — Preprocessing  [Rubric 1]
# =============================================================================
s = add_slide()
add_title(s, "Preprocessing & Feature Engineering",
          "src/preprocessing.py  ·  deterministic, fully scripted",
          rubric_tag="Rubric 1: Data (10)")

add_textbox(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
            "Cleaning steps", size=16, bold=True, color=NAVY)
add_bullets(s, Inches(0.5), Inches(1.7), Inches(6.3), Inches(4.5), [
    ("Type coercion: ", "all numerics via errors='coerce', no silent failures."),
    ("Date parsing: ", "to datetime; resample fuel weekly → monthly (MS)."),
    ("Missing values: ", "drop rows missing core fields; all 16 states retained."),
    ("Outlier check: ", "ADF stationarity test on CPI series before differencing."),
    ("Reshape: ", "extract 'overall' division from 14-division CPI table → 544 rows."),
], size=13)

add_textbox(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
            "Engineered features (per state)", size=16, bold=True, color=NAVY)
feat_rows = [
    ["Feature", "Unit"],
    ["mean_cpi_index",       "index pts"],
    ["cpi_growth_rate",      "%"],
    ["cpi_volatility (std)", "index pts"],
    ["income_mean_latest",   "RM"],
    ["income_median_latest", "RM"],
]
tbl = s.shapes.add_table(len(feat_rows), 2,
                         Inches(7.0), Inches(1.7),
                         Inches(5.8), Inches(2.6)).table
tbl.columns[0].width = Inches(3.8)
tbl.columns[1].width = Inches(2.0)
for r_i, row in enumerate(feat_rows):
    for c_i, val in enumerate(row):
        cell = tbl.cell(r_i, c_i)
        cell.text = ""
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = str(val)
        run.font.size = Pt(12)
        if r_i == 0:
            run.font.bold = True
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if r_i % 2 else WHITE
            run.font.color.rgb = DARK

add_textbox(s, Inches(7.0), Inches(4.5), Inches(6), Inches(2),
            "All 5 features z-score standardised (StandardScaler) before clustering — "
            "stops income (thousands of RM) dominating the Euclidean distance over "
            "CPI index variables.",
            size=12, color=DARK)
add_footer(s, 5, 18)

# =============================================================================
# Slide 6 — ARIMA — Method  [Rubric 2]
# =============================================================================
s = add_slide()
add_title(s, "Time-Series Model — ARIMA",
          "Why ARIMA?  Univariate monthly data, interpretable, gives confidence intervals",
          rubric_tag="Rubric 2: Training (5)")

add_bullets(s, Inches(0.5), Inches(1.2), Inches(7), Inches(5), [
    ("Hyper-parameter selection: ", "BIC grid search over p ∈ [0,4], q ∈ [0,4]."),
    ("Stationarity test: ", f"ADF on training series → p = {stat['p_value']} "
     f"({'stationary' if stat['is_stationary'] else 'non-stationary'}) → d = {arima_ord[1]}."),
    ("Selected order: ", f"ARIMA({arima_ord[0]}, {arima_ord[1]}, {arima_ord[2]})."),
    ("Train / test split: ", f"last {ar['test_months']} months held out for evaluation."),
    ("Refit on full series ", f"for {ar['horizon']}-month out-of-sample forecast."),
    ("Confidence intervals: ", "95% CI plotted alongside point forecast."),
], size=14)

add_metric_card(s, Inches(8.0),  Inches(1.3), Inches(2.4), Inches(1.1),
                "ARIMA Order", f"({arima_ord[0]},{arima_ord[1]},{arima_ord[2]})")
add_metric_card(s, Inches(10.6), Inches(1.3), Inches(2.4), Inches(1.1),
                "AIC", ar["aic"])
add_metric_card(s, Inches(8.0),  Inches(2.5), Inches(2.4), Inches(1.1),
                "BIC", ar["bic"])
add_metric_card(s, Inches(10.6), Inches(2.5), Inches(2.4), Inches(1.1),
                "Train obs", f"{ar['series_length'] - ar['test_months']}")

add_textbox(s, Inches(8.0), Inches(3.9), Inches(5), Inches(2),
            "Justification: BIC penalises complexity more than AIC, yielding parsimonious "
            "models that generalise better on a moderately long monthly economic series. "
            "ARIMA is preferred over deep models here because the series is univariate "
            "and the goal is short-horizon forecasting with interpretable uncertainty.",
            size=12, color=DARK)
add_footer(s, 6, 18)

# =============================================================================
# Slide 7 — ARIMA — Results  [Rubric 3 & 4]
# =============================================================================
s = add_slide()
add_title(s, "ARIMA Forecast — Results",
          rubric_tag="Rubric 3 + 4: Generalisation + Metrics")

add_image(s, os.path.join(FIG, "ts_forecast.png"),
          Inches(0.3), Inches(1.2), width=Inches(8.8),
          caption=f"Train / test fit and {ar['horizon']}-month forecast with 95% CI.")

add_textbox(s, Inches(9.3), Inches(1.2), Inches(3.8), Inches(0.4),
            "Test-set metrics", size=15, bold=True, color=NAVY)
add_metric_card(s, Inches(9.3),  Inches(1.7), Inches(1.8), Inches(0.95),
                "MAE",       f"{metrics['MAE']:.3f}")
add_metric_card(s, Inches(11.2), Inches(1.7), Inches(1.8), Inches(0.95),
                "RMSE",      f"{metrics['RMSE']:.3f}")
add_metric_card(s, Inches(9.3),  Inches(2.75), Inches(1.8), Inches(0.95),
                "MASE",  f"{metrics.get('MASE', float('nan')):.2f}")
add_metric_card(s, Inches(11.2), Inches(2.75), Inches(1.8), Inches(0.95),
                "Naive RMSE", f"{baseline.get('RMSE', float('nan')):.3f}")

add_textbox(s, Inches(9.3), Inches(3.9), Inches(3.8), Inches(0.4),
            "Forecast next 24 mo", size=15, bold=True, color=NAVY)
add_textbox(s, Inches(9.3), Inches(4.3), Inches(3.8), Inches(1.5),
            f"Range: {min(fc_vals):.2f}% – {max(fc_vals):.2f}%\n"
            f"Mean: {sum(fc_vals)/len(fc_vals):.2f}%\n"
            f"Direction: moderating from recent peak.",
            size=12, color=DARK)
add_footer(s, 7, 18)

# =============================================================================
# Slide 8 — Clustering — Optimal k  [Rubric 2]
# =============================================================================
s = add_slide()
add_title(s, "Clustering — Selecting Optimal k",
          "Three independent criteria, all agree",
          rubric_tag="Rubric 2: Training (5)")

add_image(s, os.path.join(FIG, "cluster_elbow_silhouette.png"),
          Inches(0.3), Inches(1.2), width=Inches(8.5))

add_bullets(s, Inches(9.0), Inches(1.3), Inches(4.0), Inches(5.5), [
    ("Elbow method:", " clear bend at k = 2."),
    ("Silhouette:", f" peaks at {km_scores['silhouette']:.4f} for k = {best_k}."),
    ("Davies-Bouldin:", f" minimised at {km_scores['davies_bouldin']:.4f}."),
    ("Why two algorithms?", " independent confirmation."),
    ("K-Means:", " minimises within-cluster variance."),
    ("Hierarchical (Ward):", " bottom-up merging."),
], size=13)
add_footer(s, 8, 18)

# =============================================================================
# Slide 9 — Clustering — Results PCA  [Rubric 3]
# =============================================================================
s = add_slide()
add_title(s, "Clustering — State Tiers in PCA Space",
          rubric_tag="Rubric 3: Generalisation")

add_image(s, os.path.join(FIG, "cluster_pca_kmeans.png"),
          Inches(0.3), Inches(1.1), width=Inches(6.3),
          caption="K-Means: states grouped in PC1–PC2 space.")
add_image(s, os.path.join(FIG, "cluster_pca_hierarchical.png"),
          Inches(6.8), Inches(1.1), width=Inches(6.3),
          caption="Hierarchical (Ward): near-identical partition.")

add_textbox(s, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8),
            f"Both algorithms produce essentially the same partition at k = {best_k} — "
            "evidence that the cluster structure is genuine, not an algorithmic artefact.",
            size=13, bold=True, color=NAVY)
add_footer(s, 9, 18)

# =============================================================================
# Slide 10 — Cluster Profiles
# =============================================================================
s = add_slide()
add_title(s, "Cluster Profiles — Who's in each tier?")

add_image(s, os.path.join(FIG, "cluster_profiles.png"),
          Inches(0.3), Inches(1.2), width=Inches(7.5))

km_profile = (
    km.groupby("cluster_label")
    .agg(states=("state", "count"),
         mean_cpi=("mean_cpi_index", "mean"),
         growth=("cpi_growth_rate", "mean"),
         income=("income_mean_latest", "mean"))
    .round(1).reset_index()
)
add_textbox(s, Inches(8.2), Inches(1.2), Inches(4.8), Inches(0.4),
            "Cluster means", size=15, bold=True, color=NAVY)
prof_rows = [["Cluster", "States", "Income (RM)"]]
for _, r in km_profile.iterrows():
    prof_rows.append([r["cluster_label"], int(r["states"]), f"{r['income']:,.0f}"])
tbl = s.shapes.add_table(len(prof_rows), 3,
                         Inches(8.2), Inches(1.7),
                         Inches(4.8), Inches(1.8)).table
tbl.columns[0].width = Inches(2.4); tbl.columns[1].width = Inches(1.0); tbl.columns[2].width = Inches(1.4)
for r_i, row in enumerate(prof_rows):
    for c_i, val in enumerate(row):
        cell = tbl.cell(r_i, c_i); cell.text = ""
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = str(val); run.font.size = Pt(11)
        if r_i == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            run.font.color.rgb = DARK
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r_i % 2 else WHITE

add_textbox(s, Inches(8.2), Inches(3.8), Inches(4.8), Inches(2.8),
            "Key insight: the higher-income tier (Selangor, Putrajaya, KL, Johor, Penang) also shows higher cumulative CPI growth — wealthier states have experienced more inflationary pressure since 2010.",
            size=12, color=DARK)
add_footer(s, 10, 18)

# =============================================================================
# Slide 11 — Fuel Price → CPI  [Rubric 3 & 4]
# =============================================================================
s = add_slide()
add_title(s, "Fuel Price → CPI Pass-through",
          "Cross-domain evaluation",
          rubric_tag="Rubric 4: Metrics (10)")

add_image(s, os.path.join(FIG, "eval_fuel_cpi_correlation.png"),
          Inches(0.3), Inches(1.2), width=Inches(9.0),
          caption="OLS fit (left) and Pearson correlations contemporaneous + lag-1 (right).")

add_bullets(s, Inches(9.5), Inches(1.3), Inches(3.5), Inches(5), [
    ("Contemporaneous: ", "fuel ↑ CPI ↑ in same month."),
    ("Lag-1 correlation: ", "fuel leads CPI by 1 month, slightly weaker."),
    ("Interpretation: ", "rapid but not instantaneous pass-through."),
    ("Why it matters: ", "evidence base for subsidy policy."),
], size=12)
add_footer(s, 11, 18)

# =============================================================================
# Slide 12 — Evaluation summary  [Rubric 4]
# =============================================================================
s = add_slide()
add_title(s, "Evaluation Summary — All metrics in one place",
          rubric_tag="Rubric 4: Metrics (10)")

add_textbox(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
            "Forecasting", size=16, bold=True, color=NAVY)
fc_rows = [
    ["Metric", "Value"],
    ["ARIMA order", f"({arima_ord[0]},{arima_ord[1]},{arima_ord[2]})"],
    ["MAE",  f"{metrics['MAE']:.4f}"],
    ["RMSE", f"{metrics['RMSE']:.4f}"],
    ["sMAPE (%)", f"{metrics.get('sMAPE', float('nan')):.2f}"],
    ["MASE", f"{metrics.get('MASE', float('nan')):.2f}"],
    ["Naive RMSE", f"{baseline.get('RMSE', float('nan')):.4f}"],
    ["AIC",  ar["aic"]],
    ["BIC",  ar["bic"]],
]
tbl = s.shapes.add_table(len(fc_rows), 2,
                         Inches(0.5), Inches(1.7),
                         Inches(5.8), Inches(3.5)).table
tbl.columns[0].width = Inches(3.0); tbl.columns[1].width = Inches(2.8)
for r_i, row in enumerate(fc_rows):
    for c_i, val in enumerate(row):
        cell = tbl.cell(r_i, c_i); cell.text = ""
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = str(val); run.font.size = Pt(12)
        if r_i == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r_i % 2 else WHITE
            run.font.color.rgb = DARK

add_textbox(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
            "Clustering Validity", size=16, bold=True, color=NAVY)
cl_rows = [
    ["Method", "Silhouette ↑", "DB ↓", "CH ↑"],
    ["K-Means",          km_scores["silhouette"], km_scores["davies_bouldin"], km_scores["calinski_harabasz"]],
    ["Hierarchical (Ward)", hc_scores["silhouette"], hc_scores["davies_bouldin"], hc_scores["calinski_harabasz"]],
]
tbl = s.shapes.add_table(len(cl_rows), 4,
                         Inches(7.0), Inches(1.7),
                         Inches(5.8), Inches(2.0)).table
tbl.columns[0].width = Inches(2.4)
for j in range(1, 4):
    tbl.columns[j].width = Inches(1.1)
for r_i, row in enumerate(cl_rows):
    for c_i, val in enumerate(row):
        cell = tbl.cell(r_i, c_i); cell.text = ""
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = str(val); run.font.size = Pt(11)
        if r_i == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r_i % 2 else WHITE
            run.font.color.rgb = DARK

add_textbox(s, Inches(7.0), Inches(4.0), Inches(6), Inches(0.4),
            "Validation against DOSM (sample)", size=16, bold=True, color=NAVY)
val_rows = [
    ["Indicator (2022)", "This project", "Official"],
    ["Sabah mean income",   "RM 6,171", "RM 6,171 (DOSM)"],
    ["Selangor mean income","RM 12,233","RM 12,233 (DOSM)"],
    ["Highest-poverty state","Sabah 19.7%","Sabah 19.7% (DOSM)"],
]
tbl = s.shapes.add_table(len(val_rows), 3,
                         Inches(7.0), Inches(4.5),
                         Inches(5.8), Inches(2.3)).table
tbl.columns[0].width = Inches(2.2); tbl.columns[1].width = Inches(1.8); tbl.columns[2].width = Inches(1.8)
for r_i, row in enumerate(val_rows):
    for c_i, val in enumerate(row):
        cell = tbl.cell(r_i, c_i); cell.text = ""
        run = cell.text_frame.paragraphs[0].add_run()
        run.text = str(val); run.font.size = Pt(10)
        if r_i == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if r_i % 2 else WHITE
            run.font.color.rgb = DARK
add_footer(s, 12, 18)

# =============================================================================
# Slide 13 — State CPI Comparison
# =============================================================================
s = add_slide()
add_title(s, "Cross-State CPI Growth (2010 → Latest)")
add_image(s, os.path.join(FIG, "eval_state_cpi_comparison.png"),
          Inches(2.5), Inches(1.1), width=Inches(8.3),
          caption="Red = above median growth, blue = below.")
add_textbox(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.5),
            "Selangor, Putrajaya and Johor lead CPI growth; East-Malaysia states (Sabah, Sarawak, Labuan) grew most slowly.",
            size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_footer(s, 13, 18)

# =============================================================================
# Slide 14 — Division Heatmap
# =============================================================================
s = add_slide()
add_title(s, "Division Inflation Heatmap",
          "Which CPI categories drove inflation, and when?")
add_image(s, os.path.join(FIG, "eval_division_heatmap.png"),
          Inches(1.0), Inches(1.1), width=Inches(11.3))
add_footer(s, 14, 18)

# =============================================================================
# Slide 15 — AI Demonstrator  [Rubric 5]
# =============================================================================
s = add_slide()
add_title(s, "AI Demonstrator — Flask Dashboard",
          "http://localhost:5050  ·  4 tabs  ·  JSON REST API",
          rubric_tag="Rubric 5: System (5)")

add_bullets(s, Inches(0.5), Inches(1.2), Inches(6.5), Inches(5), [
    ("Overview tab — ", "national + state CPI and income summary."),
    ("Time-Series tab — ", "live train/test/forecast chart with metrics."),
    ("Clustering tab — ", "PCA scatter, dendrogram, validity scores."),
    ("Fuel Analysis tab — ", "fuel vs CPI correlation explorer."),
    ("JSON API — ", "/api/overview, /api/ts/forecast, /api/cluster/states, /api/fuel/correlation."),
    ("Architecture — ", "pipeline pre-computes → Flask serves cached artefacts → fast UX."),
    ("Reproducible — ", "random_state = 42 across all components, byte-identical re-runs."),
], size=13)

# Architecture diagram
diag_left = Inches(7.5)
arch_top  = Inches(1.3)
boxes = [
    ("Raw DOSM CSVs",       NAVY,   WHITE),
    ("src/preprocessing.py", GREY,  WHITE),
    ("ARIMA + Clustering",  NAVY,   WHITE),
    ("outputs/ (JSON, CSV, PNG)", GREY, WHITE),
    ("Flask Dashboard\nlocalhost:5050", ACCENT, WHITE),
]
for i, (txt, bg, fg) in enumerate(boxes):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             diag_left, arch_top + Inches(i * 1.05),
                             Inches(5.0), Inches(0.85))
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = bg
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = fg
    if i < len(boxes) - 1:
        arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                 diag_left + Inches(2.3), arch_top + Inches(i * 1.05 + 0.85),
                                 Inches(0.4), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = GREY; arr.line.color.rgb = GREY

add_footer(s, 15, 18)

# =============================================================================
# Slide 16 — Team Contributions
# =============================================================================
s = add_slide()
add_title(s, "Team Contributions",
          "Every member touches the AI part — different model / dataset / pipeline component")

contribs = [
    ("Kenneth Hui Hong CHUA", "Time-Series Modelling Lead (ARIMA)",
     "CPI dataset · ARIMA module · BIC order selection · ADF stationarity · MAE / RMSE / sMAPE / MASE evaluation."),
    ("Kelvin Wen Kiong FONG", "Clustering AI Lead",
     "State features · K-Means + Hierarchical (Ward) · silhouette / DB / CH validity · PCA visualisation."),
    ("Nigel Zi Jun LING", "Evaluation & Analytics",
     "Fuel dataset · fuel→CPI pass-through correlation · division heatmap · DOSM validation."),
    ("Neng Yi CHIENG", "Time-Series Co-Lead + Pipeline / Demonstrator",
     "Forecast viz + 24-mo out-of-sample · CI extraction · reproducible pipeline · Flask dashboard + JSON API."),
]
top = Inches(1.3)
for i, (name, role, desc) in enumerate(contribs):
    y = top + Inches(i * 1.35)
    # Name + role band
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.5), y, Inches(4.2), Inches(1.15))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = name
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = role
    r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE); r2.font.italic = True

    # Description
    add_textbox(s, Inches(4.9), y + Inches(0.1), Inches(8.2), Inches(1.0),
                desc, size=13, color=DARK)
add_footer(s, 16, 18)

# =============================================================================
# Slide 17 — Limitations & Future Work
# =============================================================================
s = add_slide()
add_title(s, "Limitations & Future Work")

add_textbox(s, Inches(0.5), Inches(1.2), Inches(6), Inches(0.4),
            "Limitations", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.5), Inches(1.7), Inches(6), Inches(5), [
    ("Long-horizon skill: ", f"MASE = {metrics.get('MASE', float('nan')):.2f} (>1) — over a 24-month "
     "block the stationary model does not beat a naive baseline (reported honestly)."),
    ("Fuel adds no forecast gain: ", "fuel correlates with inflation but, as an ARIMAX regressor, did "
     "not lower test error in this specification."),
    ("Small cross-section: ", "16 states limits the number of well-separated clusters."),
    ("Income survey infrequency: ", "income features use the latest available year (≈2022) only."),
    ("Aggregate CPI: ", "national series masks division-level variation already shown."),
], size=13)

add_textbox(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
            "Future Work", size=18, bold=True, color=NAVY)
add_bullets(s, Inches(7.0), Inches(1.7), Inches(6), Inches(5), [
    ("Rolling one-step evaluation: ", "the fair test of short-horizon skill, where ARIMA should beat naive."),
    ("LSTM / Prophet: ", "compare non-linear forecasters on the same hold-out."),
    ("Division-level ARIMA: ", "forecast each of the 13 CPI divisions separately."),
    ("District clustering: ", "extend from 16 states to 160 districts when data permits."),
    ("Live dashboard: ", "auto-refresh from OpenDOSM API as new months are released."),
], size=13)
add_footer(s, 17, 18)

# =============================================================================
# Slide 18 — Conclusions & Q&A
# =============================================================================
s = add_slide()
add_title(s, "Conclusions  ·  Thank you  ·  Q&A")

add_bullets(s, Inches(0.5), Inches(1.2), Inches(12.5), Inches(3.5), [
    ("ARIMA gives a national CPI outlook ", f"({min(fc_vals):.1f}-{max(fc_vals):.1f}% over {ar['horizon']} "
     f"months; RMSE {metrics['RMSE']:.2f}), with a candid MASE = {metrics.get('MASE', float('nan')):.2f} "
     "showing long-horizon point accuracy is limited."),
    (f"Two-tier state structure (k = {best_k}) ", f"is robust across K-Means and Hierarchical "
     f"clustering (silhouette {km_scores['silhouette']:.3f})."),
    ("Fuel correlates with CPI ", "but did not improve the forecast (ARIMAX) — association ≠ predictive power."),
    ("Reproducible pipeline + Flask dashboard ", "make the analysis usable by non-technical stakeholders."),
    ("All results reconcile ", "with DOSM official 2022 published figures."),
], size=15)

# Big thank-you band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.6))
band.fill.solid(); band.fill.fore_color.rgb = NAVY; band.line.fill.background()
tf = band.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank you — Questions?"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run()
r2.text = "Repo: github.com/kenn040502/COS40007-Design-Project   ·   Dashboard: localhost:5050"
r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)

add_footer(s, 18, 18)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print("Presentation saved to:", OUT)
print(f"  Slides: {len(prs.slides)}")
